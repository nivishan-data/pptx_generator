from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO

# Generate a PowerPoint presentation
def generate_ppt(slide_content, images, file_name="generated_presentation.pptx"):
    prs = Presentation()

    for i, slide in enumerate(slide_content):
        # Use a simple layout for the slide
        slide_layout = prs.slide_layouts[1]
        slide_object = prs.slides.add_slide(slide_layout)
        title = slide_object.shapes.title
        content = slide_object.placeholders[1]

        # Set the title
        title.text = slide['title']

        # Add formatted content
        tf = content.text_frame
        for line in slide['content'].split('\n'):
            if line.startswith("*"):  # Bullet points
                p = tf.add_paragraph()
                p.text = line[1:].strip()
                p.level = 1
            elif line.startswith("**"):  # Bold text
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = line[2:].strip()
                run.bold = True
                p.alignment = PP_ALIGN.LEFT
            else:  # Regular text
                p = tf.add_paragraph()
                p.text = line.strip()
                p.alignment = PP_ALIGN.LEFT

        # Add image to the slide if available
        if images and i < len(images):
            image_stream = BytesIO()
            images[i].save(image_stream, format='PNG')
            image_stream.seek(0)
            
            # Place the image at a specific location to avoid overlap with text
            slide_object.shapes.add_picture(image_stream, Inches(5), Inches(3), width=Inches(4), height=Inches(3))

    prs.save(file_name)
    return file_name
